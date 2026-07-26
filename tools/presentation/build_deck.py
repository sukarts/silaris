"""Génère la présentation commerciale SILARIS (PowerPoint 16:9).

Les visuels proviennent des captures réelles de l'application (apps/web/captures),
produites par le spec Playwright `capture.spec.ts` sur le jeu de démonstration.
"""
import os
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Emu, Inches, Pt

ROOT = Path(__file__).resolve().parents[2]
CAPTURES = Path(os.environ.get("CAPTURE_DIR", ROOT / "apps/web/captures"))
EXTRA = Path(__file__).resolve().parent
OUT = Path(os.environ.get("DECK_OUT", ROOT / "SILARIS-presentation.pptx"))

# Palette reprise des tokens de l'application.
INK = RGBColor(0x1F, 0x24, 0x30)
INK2 = RGBColor(0x4B, 0x53, 0x63)
INK3 = RGBColor(0x8A, 0x93, 0xA3)
ACCENT = RGBColor(0xE8, 0x66, 0x3D)
SEA = RGBColor(0x3D, 0x7F, 0xA6)
OK = RGBColor(0x2E, 0x9E, 0x6B)
WARN = RGBColor(0xC9, 0x8A, 0x21)
PAPER = RGBColor(0xF7, 0xF8, 0xFA)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LINE = RGBColor(0xD9, 0xDC, 0xE2)

W, H = Inches(13.333), Inches(7.5)

prs = Presentation()
prs.slide_width, prs.slide_height = W, H
BLANK = prs.slide_layouts[6]


# ---------------------------------------------------------------- primitives


def slide(bg=WHITE):
    s = prs.slides.add_slide(BLANK)
    rect = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, W, H)
    rect.fill.solid()
    rect.fill.fore_color.rgb = bg
    rect.line.fill.background()
    rect.shadow.inherit = False
    return s


def text(s, txt, x, y, w, h, size=18, color=INK, bold=False, align=PP_ALIGN.LEFT,
         font="Helvetica Neue", spacing=1.0, anchor=MSO_ANCHOR.TOP):
    box = s.shapes.add_textbox(x, y, w, h)
    tf = box.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    lines = txt.split("\n")
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.line_spacing = spacing
        run = p.add_run()
        run.text = line
        run.font.size = Pt(size)
        run.font.color.rgb = color
        run.font.bold = bold
        run.font.name = font
    return box


def chip(s, txt, x, y, color, w=None, size=11):
    w = w or Inches(1.25)
    h = Inches(0.3)
    box = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
    box.fill.solid()
    box.fill.fore_color.rgb = color
    box.line.fill.background()
    box.shadow.inherit = False
    box.adjustments[0] = 0.5
    tf = box.text_frame
    tf.margin_left = tf.margin_right = 0
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = txt
    run.font.size = Pt(size)
    run.font.bold = True
    run.font.color.rgb = WHITE
    run.font.name = "Helvetica Neue"
    return box


def card(s, x, y, w, h, fill=PAPER, border=None):
    box = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
    box.fill.solid()
    box.fill.fore_color.rgb = fill
    box.adjustments[0] = 0.04
    box.shadow.inherit = False
    if border:
        box.line.color.rgb = border
        box.line.width = Pt(1)
    else:
        box.line.fill.background()
    return box


def wordmark(s, x, y, size=20):
    """SILA + RIS en accent, comme dans l'application."""
    box = s.shapes.add_textbox(x, y, Inches(3), Inches(0.5))
    tf = box.text_frame
    tf.margin_left = tf.margin_top = 0
    p = tf.paragraphs[0]
    for part, col in (("SILA", INK), ("RIS", ACCENT)):
        r = p.add_run()
        r.text = part
        r.font.size = Pt(size)
        r.font.bold = True
        r.font.color.rgb = col
        r.font.name = "Helvetica Neue"
    p.runs[0].font._rPr.set("spc", "180")
    return box


def shot(s, name, x, y, w, shadow=True):
    """Capture d'écran cadrée, ratio préservé (les captures font 1500x940)."""
    path = CAPTURES / name if (CAPTURES / name).exists() else EXTRA / name
    pic = s.shapes.add_picture(str(path), x, y, width=w)
    border = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, pic.height)
    border.fill.background()
    border.line.color.rgb = LINE
    border.line.width = Pt(1.25)
    border.adjustments[0] = 0.02
    border.shadow.inherit = shadow
    return pic


def section_header(s, eyebrow, title, sub=None):
    text(s, eyebrow.upper(), Inches(0.85), Inches(0.55), Inches(8), Inches(0.3),
         size=11, color=ACCENT, bold=True)
    text(s, title, Inches(0.85), Inches(0.85), Inches(11.6), Inches(0.6), size=30, bold=True)
    if sub:
        text(s, sub, Inches(0.85), Inches(1.5), Inches(11.6), Inches(0.4), size=14, color=INK2)


def footer(s, n):
    text(s, "SILARIS", Inches(0.85), Inches(6.95), Inches(2), Inches(0.3),
         size=9, color=INK3, bold=True)
    text(s, str(n), Inches(12), Inches(6.95), Inches(0.5), Inches(0.3),
         size=9, color=INK3, align=PP_ALIGN.RIGHT)


# ------------------------------------------------------------------- 1. Cover
s = slide(INK)
band = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(0.18), H)
band.fill.solid()
band.fill.fore_color.rgb = ACCENT
band.line.fill.background()
band.shadow.inherit = False

box = s.shapes.add_textbox(Inches(1.1), Inches(2.0), Inches(9), Inches(1.2))
p = box.text_frame.paragraphs[0]
for part, col in (("SILA", WHITE), ("RIS", ACCENT)):
    r = p.add_run()
    r.text = part
    r.font.size = Pt(66)
    r.font.bold = True
    r.font.color.rgb = col
    r.font.name = "Helvetica Neue"

text(s, "Le système d'information du transitaire", Inches(1.15), Inches(3.3), Inches(10), Inches(0.6),
     size=27, color=WHITE)
text(s, "Maritime · Aérien · Routier · Douane · Facturation · Portail client",
     Inches(1.15), Inches(4.05), Inches(10), Inches(0.4), size=15, color=INK3)

for i, (val, lab) in enumerate([("19", "modules"), ("100", "tables"), ("169", "routes API"), ("109", "tests")]):
    x = Inches(1.15 + i * 2.1)
    text(s, val, x, Inches(4.95), Inches(1.9), Inches(0.6), size=34, bold=True, color=ACCENT)
    text(s, lab, x, Inches(5.6), Inches(1.9), Inches(0.3), size=12, color=INK3)

text(s, "Présentation de la solution — MVP en production", Inches(1.15), Inches(6.5), Inches(10), Inches(0.4),
     size=12, color=INK3)

# --------------------------------------------------------------- 2. Le constat
s = slide()
section_header(s, "Le constat", "Le transit se pilote encore au tableur et au téléphone")
items = [
    ("Information éclatée", "Le dossier vit dans un classeur, la position du conteneur dans un mail, la facture dans un autre logiciel. Personne n'a la vue complète."),
    ("Le client rappelle", "« Où est ma marchandise ? » — chaque appel mobilise un exploitant qui doit rouvrir trois outils pour répondre."),
    ("Rien n'est traçable", "Qui a modifié l'ETA ? Qui a validé la facture ? Sans piste d'audit, la réponse se perd dans les souvenirs."),
    ("La preuve se perd", "Le bon de livraison signé repart en papier, se froisse dans un camion, et manque le jour du litige."),
]
for i, (t, d) in enumerate(items):
    x = Inches(0.85 + (i % 2) * 6.05)
    y = Inches(2.1 + (i // 2) * 2.25)
    card(s, x, y, Inches(5.7), Inches(1.95))
    text(s, t, x + Inches(0.4), y + Inches(0.35), Inches(5), Inches(0.35), size=17, bold=True)
    text(s, d, x + Inches(0.4), y + Inches(0.85), Inches(5), Inches(1), size=12.5, color=INK2, spacing=1.25)
footer(s, 2)

# --------------------------------------------------------------- 3. La réponse
s = slide(INK)
text(s, "LA RÉPONSE", Inches(0.85), Inches(0.55), Inches(8), Inches(0.3), size=11, color=ACCENT, bold=True)
text(s, "Un seul dossier, du booking à l'encaissement",
     Inches(0.85), Inches(0.95), Inches(11.6), Inches(0.6), size=30, bold=True, color=WHITE)
text(s, "Chaque expédition porte son mode de transport, ses conteneurs, ses documents, ses coûts et ses factures.\nLe client suit son dossier lui-même. Le suivi maritime s'actualise seul.",
     Inches(0.85), Inches(1.75), Inches(11.3), Inches(1), size=14.5, color=INK3, spacing=1.35)

pillars = [
    ("Un dossier unique", "Maritime, aérien et routier dans\nla même chaîne, du devis à la facture.", ACCENT),
    ("Le client autonome", "Portail dédié et page de suivi publique\npar numéro de conteneur ou de BL.", SEA),
    ("Le suivi automatique", "Les positions conteneurs remontent\nseules depuis onze compagnies.", OK),
    ("Chaque geste tracé", "Piste d'audit, rôles fins et cloisonnement\nétanche entre sociétés.", WARN),
]
for i, (t, d, col) in enumerate(pillars):
    x = Inches(0.85 + i * 3.05)
    card(s, x, Inches(3.15), Inches(2.8), Inches(2.6), fill=RGBColor(0x2A, 0x30, 0x3D))
    bar = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x + Inches(0.35), Inches(3.5), Inches(0.5), Inches(0.08))
    bar.fill.solid()
    bar.fill.fore_color.rgb = col
    bar.line.fill.background()
    bar.shadow.inherit = False
    text(s, t, x + Inches(0.35), Inches(3.8), Inches(2.2), Inches(0.6), size=16, bold=True, color=WHITE)
    text(s, d, x + Inches(0.35), Inches(4.5), Inches(2.2), Inches(1.1), size=11.5, color=INK3, spacing=1.3)
footer(s, 3)

# ------------------------------------------------- 4. Cartographie des modules
s = slide()
section_header(s, "Périmètre", "Vingt-quatre briques fonctionnelles, dix-neuf déjà livrées",
               "Le décompte porte sur les briques visibles à l'usage.")
groups = [
    ("Opérations", [("Dossiers", 1), ("Bookings", 1), ("Conteneurs", 1), ("Aérien", 1), ("Routier", 1), ("Télématique", 1)]),
    ("Commercial", [("CRM", 1), ("Cotations", 1), ("Facturation", 1), ("Recouvrement", 0)]),
    ("Ressources", [("Documents", 1), ("Référentiels", 1), ("Reporting", 1), ("Entrepôt / stock", 0)]),
    ("Plateforme", [("Identité & rôles", 1), ("Multi-société", 1), ("Notifications", 1), ("Audit", 1), ("Recherche", 1),
                    ("Suivi maritime", 1), ("Portail client", 1), ("Odoo", 2), ("Douane EDI", 0), ("Mobile chauffeur", 0)]),
]
colors = {1: OK, 2: WARN, 0: INK3}
labels = {1: "Livré", 2: "Prêt", 0: "À venir"}
y = Inches(2.02)
RIGHT = Inches(12.9)
for title, mods in groups:
    text(s, title, Inches(0.85), y, Inches(2), Inches(0.35), size=15, bold=True, color=INK)
    x = Inches(2.75)
    rows = 1
    for name, state in mods:
        w = Inches(0.22 + 0.105 * len(name))
        # Les groupes longs (Plateforme) repassent à la ligne plutôt que de
        # déborder de la diapositive.
        if x + w > RIGHT:
            x = Inches(2.75)
            y += Inches(0.55)
            rows += 1
        card(s, x, y - Inches(0.03), w, Inches(0.42), fill=WHITE, border=colors[state])
        text(s, name, x, y + Inches(0.045), w, Inches(0.3), size=11.5,
             color=colors[state] if state != 1 else INK, bold=True, align=PP_ALIGN.CENTER)
        x += w + Inches(0.14)
    y += Inches(0.95)

lx = Inches(0.85)
for state in (1, 2, 0):
    dot = s.shapes.add_shape(MSO_SHAPE.OVAL, lx, Inches(6.62), Inches(0.13), Inches(0.13))
    dot.fill.solid()
    dot.fill.fore_color.rgb = colors[state]
    dot.line.fill.background()
    dot.shadow.inherit = False
    cap = {1: "Livré et en production", 2: "Développé, activation sur demande", 0: "Feuille de route"}[state]
    text(s, cap, lx + Inches(0.25), Inches(6.59), Inches(3.6), Inches(0.3), size=11, color=INK2)
    lx += Inches(3.9)
footer(s, 4)

# ------------------------------------------------------- Fiches module (shots)
MODULES = [
    ("Pilotage", "Tableau de bord", "01-tableau-de-bord.png",
     "Dossiers actifs, conteneurs en cours, retards détectés et chiffre d'affaires du mois en un écran.",
     ["Alertes de retard calculées sur l'écart à l'ETA initiale", "Volumes import / export sur six mois",
      "Accès direct aux derniers dossiers", "Recherche globale ⌘K sur dossiers, clients et conteneurs"]),
    ("Opérations", "Dossiers", "02-dossiers.png",
     "Le cœur du métier : chaque expédition, quel que soit son mode, suit un workflow d'étapes configurable.",
     ["Import, export et transit — maritime FCL/LCL, aérien, routier", "Référence générée au format choisi par le transitaire",
      "Filtres par statut, mode, client et agence", "Retards signalés automatiquement"]),
    ("Opérations", "Détail du dossier", "03-dossier-detail.png",
     "Toute la vie de l'expédition sur une page : étapes, timeline, documents, conteneurs et coûts.",
     ["Stepper de workflow avec conditions de passage", "Timeline horodatée de chaque événement",
      "Documents attendus et manquants signalés", "Marge calculée en continu"]),
    ("Opérations", "Bookings", "04-bookings.png",
     "La réservation auprès de la compagnie, du numéro de booking jusqu'au connaissement.",
     ["Compagnie maritime en autocomplétion", "Navire, voyage, ports et dates",
      "Rattachement automatique au dossier", "Numéros BL, HBL et MBL"]),
    ("Opérations", "Conteneurs", "05-conteneurs.png",
     "Le parc conteneur et son affectation aux dossiers, avec scellés et suivi de position.",
     ["Types ISO et caractéristiques techniques", "Affectation multi-dossiers avec scellé",
      "Position remontée par la compagnie", "Détention et surestaries suivies"]),
    ("Opérations", "Aérien", "06-aerien.png",
     "Le fret aérien avec ses propres codes : LTA, compagnies, aéroports et poids taxable.",
     ["Master et House Air Waybill", "Aéroports IATA en autocomplétion",
      "Poids brut, volumétrique et taxable", "Vols et escales"]),
    ("Opérations", "Routier — missions", "07-routier-missions.png",
     "Le pré et post-acheminement : missions, étapes, preuve de livraison signée.",
     ["Livraison, enlèvement et transfert", "Étapes géolocalisées avec détection d'arrivée",
      "Sous-traitance : transporteur affrété identifié", "Bon de livraison PDF signé"]),
    ("Opérations", "Routier — flotte", "08-routier-flotte.png",
     "Camions, remorques et chauffeurs, en propre ou fournis par un prestataire.",
     ["Propriétaire distingué : flotte propre ou affrétée", "Visites techniques et assurances suivies",
      "Permis et échéances des chauffeurs", "Affectation croisée refusée par le serveur"]),
    ("Opérations", "Routier — balises GPS", "09-routier-balises.png",
     "Le dernier kilomètre suivi par balise, du port jusqu'à la livraison.",
     ["Compatible Teltonika, Queclink, Traccar, Flespi", "Clé d'ingestion par balise, jamais un compte utilisateur",
      "Géorepérage à 250 m pour détecter l'arrivée", "Position visible du client pendant la livraison"]),
    ("Commercial", "CRM", "10-crm.png",
     "Clients, prospects et fournisseurs — dont les transporteurs sous-traitants.",
     ["Personne physique ou morale, avec RCCM et pièce d'identité", "Code client généré automatiquement",
      "Secteur d'activité, pays et indicatif avec drapeaux", "Invitation du client au portail en un clic"]),
    ("Commercial", "Cotations", "11-cotations.png",
     "Le devis chiffré, envoyé au client, accepté ou refusé depuis son portail.",
     ["Lignes de coûts et de ventes, marge calculée", "Validité et relances",
      "PDF à la marque du transitaire", "Acceptation en ligne par le client"]),
    ("Commercial", "Facturation", "12-facturation.png",
     "Facture, proforma et avoir, avec numérotation au format de la société.",
     ["Numérotation configurable par société", "TVA, devises et échéances",
      "États de règlement suivis", "PDF disponible au portail client"]),
    ("Ressources", "Documents", "13-documents.png",
     "La GED du dossier : connaissements, factures fournisseurs, déclarations douanières.",
     ["Stockage objet chiffré (Cloudflare R2)", "Documents requis par étape de workflow",
      "Liens de téléchargement signés", "Partage sélectif avec le client"]),
    ("Plateforme", "Administration", "14-administration.png",
     "Utilisateurs, rôles et piste d'audit — qui a fait quoi, et quand.",
     ["Onze rôles système, 119 permissions atomiques", "Invitation par email avec mot de passe temporaire",
      "Piste d'audit horodatée et exportable", "Rattachement des utilisateurs aux agences"]),
    ("Plateforme", "Paramètres — société", "15-parametres-societe.png",
     "L'identité du transitaire : mentions légales, logo, formats de numérotation.",
     ["Logo repris sur tous les documents et le portail", "Format de référence dossier avec aperçu en direct",
      "Format de numéro de facture", "Mentions légales et fiscales"]),
    ("Plateforme", "Paramètres — agences", "16-parametres-agences.png",
     "Implantations propres et correspondants partenaires à l'étranger.",
     ["Code agence normalisé UN/LOCODE (CIABJ, BEANR)", "Correspondant externe distingué de l'agence propre",
      "Pays, ville et adresse", "Fuseaux horaires mondiaux"]),
    ("Plateforme", "Profil et sécurité", "17-profil-mfa.png",
     "Chaque utilisateur gère son compte et active sa double authentification.",
     ["TOTP compatible Google Authenticator et 1Password", "Codes de récupération à usage unique",
      "Changement de mot de passe imposé à la première connexion", "Sessions et jetons révocables"]),
    ("Client", "Suivi public", "18-suivi-public.png",
     "Sans compte ni mot de passe : un numéro de conteneur, de BL ou de LTA suffit.",
     ["Recherche par conteneur, BL, HBL, MBL ou LTA", "Étapes franchies et arrivée estimée",
      "Position du véhicule pendant la livraison", "À la marque du transitaire, pas à celle de SILARIS"]),
    ("Client", "Portail client", "19-portail-connexion.png",
     "L'espace dédié du client : ses dossiers, ses factures, ses cotations, ses bons de livraison.",
     ["Compte distinct des utilisateurs internes", "Acceptation des cotations en ligne",
      "Factures et bons de livraison en PDF", "Aucune donnée interne : ni marge, ni sous-traitant"]),
]

n = 5
for eyebrow, title, img, lead, bullets in MODULES:
    s = slide()
    text(s, eyebrow.upper(), Inches(0.85), Inches(0.5), Inches(6), Inches(0.28),
         size=10.5, color=ACCENT, bold=True)
    text(s, title, Inches(0.85), Inches(0.8), Inches(6), Inches(0.55), size=27, bold=True)
    text(s, lead, Inches(0.85), Inches(1.5), Inches(4.0), Inches(1.5), size=12.5, color=INK2, spacing=1.3)
    rule = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.85), Inches(3.05), Inches(0.6), Inches(0.035))
    rule.fill.solid()
    rule.fill.fore_color.rgb = ACCENT
    rule.line.fill.background()
    rule.shadow.inherit = False
    y = Inches(3.42)
    for b in bullets:
        dot = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(0.9), y + Inches(0.085), Inches(0.085), Inches(0.085))
        dot.fill.solid()
        dot.fill.fore_color.rgb = SEA
        dot.line.fill.background()
        dot.shadow.inherit = False
        text(s, b, Inches(1.18), y, Inches(3.7), Inches(0.6), size=11.5, color=INK, spacing=1.25)
        y += Inches(0.78)
    # La capture occupe la moitié droite, centrée verticalement sur la diapositive.
    shot(s, img, Inches(5.15), Inches(1.42), Inches(7.35))
    footer(s, n)
    n += 1

# ------------------------------------------------- Documents produits (PDF)
s = slide()
section_header(s, "Documents", "Des pièces à votre marque, pas à celle de SILARIS",
               "Factures, proformas, avoirs, cotations, étiquettes colis et bons de livraison signés.")
shot(s, "bon-livraison.png", Inches(0.85), Inches(2.2), Inches(3.3))
notes = [
    ("Bon de livraison signé", "Marchandise remise, destinataire, lieu et heure,\nsignature manuscrite capturée sur place."),
    ("Facture et proforma", "Numérotation au format de la société,\nTVA et échéances, disponible au portail."),
    ("Cotation", "Lignes chiffrées, validité, acceptation\nen ligne par le client."),
]
y = Inches(2.35)
for t, d in notes:
    card(s, Inches(4.7), y, Inches(7.8), Inches(1.32))
    text(s, t, Inches(5.05), y + Inches(0.22), Inches(7), Inches(0.35), size=16, bold=True)
    text(s, d, Inches(5.05), y + Inches(0.65), Inches(7), Inches(0.6), size=12, color=INK2, spacing=1.25)
    y += Inches(1.5)
footer(s, n)
n += 1

# ------------------------------------------------------------- Intégrations
s = slide()
section_header(s, "Intégrations", "Branché sur le monde réel, pas sur des données de démonstration")
integrations = [
    ("JSONCargo", "Suivi conteneurs", "Onze compagnies maritimes interrogées\npar numéro de conteneur ou de BL.", OK, "Actif"),
    ("Resend", "Emails transactionnels", "Domaine silaris.co vérifié : invitations,\nréinitialisations, alertes de retard.", OK, "Actif"),
    ("Meilisearch", "Recherche globale", "Dossiers, clients et conteneurs indexés,\nrecherche instantanée au ⌘K.", OK, "Actif"),
    ("Cloudflare R2", "Stockage documents", "Stockage objet compatible S3,\nliens de téléchargement signés.", OK, "Actif"),
    ("Odoo", "Comptabilité", "Synchronisation des factures et tiers\npar JSON-RPC — développé et testé.", WARN, "Prêt"),
    ("Balises GPS", "Télématique", "Teltonika, Queclink, Traccar, Flespi :\ntout émetteur HTTP est accepté.", OK, "Actif"),
]
for i, (name, cat, desc, col, state) in enumerate(integrations):
    x = Inches(0.85 + (i % 3) * 4.0)
    y = Inches(2.1 + (i // 3) * 2.4)
    card(s, x, y, Inches(3.7), Inches(2.05), fill=WHITE, border=LINE)
    text(s, name, x + Inches(0.35), y + Inches(0.3), Inches(2.3), Inches(0.35), size=17, bold=True)
    chip(s, state, x + Inches(2.65), y + Inches(0.32), col, w=Inches(0.72), size=9.5)
    text(s, cat, x + Inches(0.35), y + Inches(0.72), Inches(3), Inches(0.3), size=11, color=ACCENT, bold=True)
    text(s, desc, x + Inches(0.35), y + Inches(1.1), Inches(3.1), Inches(0.8), size=11.5, color=INK2, spacing=1.25)
footer(s, n)
n += 1

# ---------------------------------------------------------------- Sécurité
s = slide(INK)
text(s, "SÉCURITÉ & CONFORMITÉ", Inches(0.85), Inches(0.55), Inches(8), Inches(0.3),
     size=11, color=ACCENT, bold=True)
text(s, "Le cloisonnement est garanti par la base de données",
     Inches(0.85), Inches(0.95), Inches(11.6), Inches(0.6), size=28, bold=True, color=WHITE)
text(s, "Chaque société est isolée par les politiques de sécurité au niveau ligne de PostgreSQL : même en cas d'erreur applicative,\nune requête ne peut pas franchir la frontière d'un tenant.",
     Inches(0.85), Inches(1.75), Inches(11.4), Inches(0.9), size=13.5, color=INK3, spacing=1.35)
sec = [
    ("Isolation multi-société", "Row Level Security PostgreSQL sur 57 tables"),
    ("Rôles fins", "11 rôles système, 119 permissions atomiques"),
    ("Double authentification", "TOTP et codes de récupération"),
    ("Piste d'audit", "Chaque écriture horodatée et attribuée"),
    ("Confidentialité client", "Ni marge, ni sous-traitant, ni chauffeur exposés"),
    ("Vie privée des chauffeurs", "Positions bornées, précision réduite côté client"),
]
for i, (t, d) in enumerate(sec):
    x = Inches(0.85 + (i % 3) * 4.0)
    y = Inches(3.15 + (i // 3) * 1.75)
    card(s, x, y, Inches(3.7), Inches(1.45), fill=RGBColor(0x2A, 0x30, 0x3D))
    text(s, t, x + Inches(0.35), y + Inches(0.28), Inches(3.1), Inches(0.35), size=14.5, bold=True, color=WHITE)
    text(s, d, x + Inches(0.35), y + Inches(0.72), Inches(3.1), Inches(0.5), size=11.5, color=INK3, spacing=1.2)
footer(s, n)
n += 1

# ------------------------------------------------------------ Ce qui est livré
s = slide()
section_header(s, "MVP", "Ce qui tourne aujourd'hui en production",
               "Déployé sur silaris.co — accessible aux exploitants comme aux clients.")
delivered = [
    ("Chaîne opérationnelle complète", "Dossiers, bookings, conteneurs, aérien, routier — du devis à la facture."),
    ("Suivi maritime automatique", "Onze compagnies interrogées chaque jour, actualisation manuelle possible."),
    ("Portail et suivi public", "Le client consulte ses dossiers, accepte ses cotations, télécharge ses pièces."),
    ("Documents à votre marque", "Factures, cotations et bons de livraison signés, sous votre logo."),
    ("Emails réels", "Invitations, réinitialisations et alertes envoyés depuis votre domaine."),
    ("Télématique routière", "Balises GPS, géorepérage et preuve de livraison signée sur place."),
]
y = Inches(2.15)
for i, (t, d) in enumerate(delivered):
    x = Inches(0.85 + (i % 2) * 6.05)
    yy = y + Inches((i // 2) * 1.55)
    card(s, x, yy, Inches(5.7), Inches(1.3), fill=WHITE, border=OK)
    text(s, "✓", x + Inches(0.3), yy + Inches(0.26), Inches(0.4), Inches(0.4), size=17, bold=True, color=OK)
    text(s, t, x + Inches(0.78), yy + Inches(0.26), Inches(4.6), Inches(0.35), size=15, bold=True)
    text(s, d, x + Inches(0.78), yy + Inches(0.68), Inches(4.7), Inches(0.5), size=11.5, color=INK2, spacing=1.2)
footer(s, n)
n += 1

# -------------------------------------------------------------- Feuille de route
s = slide()
section_header(s, "Feuille de route", "Ce qui reste à construire")
roadmap = [
    ("Court terme", ACCENT, [
        ("Activation Odoo", "Développé et testé — attend les accès de votre instance."),
        ("Sous-domaines par client", "societe.silaris.co, résolveur déjà en place."),
        ("Préférences de notification", "Choisir par utilisateur ce qui déclenche un email."),
    ]),
    ("Moyen terme", SEA, [
        ("Recouvrement", "Relances automatiques, balance âgée, échéancier."),
        ("Entrepôt et stock", "Réception, emplacements, inventaire, sorties."),
        ("Application chauffeur", "POD signée par le chauffeur lui-même, hors connexion."),
    ]),
    ("Long terme", INK3, [
        ("Douane EDI", "Transmission dématérialisée des déclarations."),
        ("Tarifs négociés", "Grilles compagnies importées, cotation automatique."),
        ("Analytique avancée", "Rentabilité par client, par ligne, par agence."),
    ]),
]
x = Inches(0.85)
for horizon, col, items in roadmap:
    bar = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, Inches(2.1), Inches(3.85), Inches(0.09))
    bar.fill.solid()
    bar.fill.fore_color.rgb = col
    bar.line.fill.background()
    bar.shadow.inherit = False
    text(s, horizon, x, Inches(2.35), Inches(3.7), Inches(0.4), size=18, bold=True, color=col)
    y = Inches(3.0)
    for t, d in items:
        card(s, x, y, Inches(3.85), Inches(1.2), fill=PAPER)
        text(s, t, x + Inches(0.3), y + Inches(0.22), Inches(3.3), Inches(0.3), size=13.5, bold=True)
        text(s, d, x + Inches(0.3), y + Inches(0.6), Inches(3.3), Inches(0.5), size=10.5, color=INK2, spacing=1.2)
        y += Inches(1.35)
    x += Inches(4.05)
footer(s, n)
n += 1

# ------------------------------------------------------------------ Technique
s = slide()
section_header(s, "Socle technique", "Construit pour durer et pour être repris")
tech = [
    ("Backend", "Laravel 12 · PHP 8.3", "Monolithe modulaire en architecture hexagonale :\ndix-neuf modules étanches, testables un à un."),
    ("Base de données", "PostgreSQL 16", "Cent tables, isolation par Row Level Security,\nmigrations versionnées."),
    ("Frontend", "Next.js 15 · React 19", "Vingt-sept écrans, client d'API typé\ngénéré depuis la spécification OpenAPI."),
    ("Qualité", "109 tests · Pint · PHPStan", "Chaque fusion passe la suite complète,\nl'analyse statique et le formatage."),
    ("Déploiement", "Docker · Render · GHCR", "Images construites par la CI,\nmigrations jouées au pré-déploiement."),
    ("Exploitation", "Horizon · Scheduler", "Files d'attente supervisées,\ntâches planifiées et outbox transactionnel."),
]
for i, (cat, stack, desc) in enumerate(tech):
    x = Inches(0.85 + (i % 3) * 4.0)
    y = Inches(2.15 + (i // 3) * 2.35)
    card(s, x, y, Inches(3.7), Inches(2.0), fill=WHITE, border=LINE)
    text(s, cat.upper(), x + Inches(0.35), y + Inches(0.28), Inches(3), Inches(0.25), size=10, color=ACCENT, bold=True)
    text(s, stack, x + Inches(0.35), y + Inches(0.6), Inches(3.1), Inches(0.4), size=15, bold=True)
    text(s, desc, x + Inches(0.35), y + Inches(1.1), Inches(3.1), Inches(0.8), size=11, color=INK2, spacing=1.25)
footer(s, n)
n += 1

# ---------------------------------------------------------------------- Clôture
s = slide(INK)
band = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(0.18), H)
band.fill.solid()
band.fill.fore_color.rgb = ACCENT
band.line.fill.background()
band.shadow.inherit = False
text(s, "Le MVP est en production", Inches(1.15), Inches(2.4), Inches(10), Inches(0.8),
     size=40, bold=True, color=WHITE)
text(s, "Les modules opérationnels, commerciaux et clients fonctionnent aujourd'hui sur des données réelles.\nLa suite se construit au rythme de vos priorités.",
     Inches(1.15), Inches(3.5), Inches(10.5), Inches(1), size=15, color=INK3, spacing=1.4)
text(s, "app.silaris.co", Inches(1.15), Inches(5.1), Inches(5), Inches(0.4), size=17, bold=True, color=ACCENT)
text(s, "portail client · suivi public · espace exploitant", Inches(1.15), Inches(5.6), Inches(7), Inches(0.4),
     size=12, color=INK3)

prs.save(str(OUT))
print(f"{OUT} — {len(prs.slides.__iter__.__self__._sldIdLst)} diapositives")
